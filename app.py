import argparse
import base64
import logging
import os
import io
import tempfile
import time
import zipfile
import csv
import re
import html
from html.parser import HTMLParser
from pathlib import Path
from typing import Tuple

# ---------- Optional imports ----------
try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except Exception:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except Exception:
    OCR_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except Exception:
    PANDAS_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except Exception:
    OPENPYXL_AVAILABLE = False

try:
    from email import policy
    from email.parser import BytesParser
    EMAIL_AVAILABLE = True
except Exception:
    EMAIL_AVAILABLE = False

try:
    import extract_msg
    MSG_AVAILABLE = True
except Exception:
    MSG_AVAILABLE = False

# ---------- Logging ----------
logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(message)s')
logger = logging.getLogger("extractor")
logging.getLogger("extract_msg").setLevel(logging.WARNING)
OCR_LANG = os.environ.get("OCR_LANG", "eng")

# ---------- Utilities ----------
def _to_str(x):
    if isinstance(x, bytes):
        return x.decode("utf-8", errors="replace")
    return str(x or "")

class _HTMLToText(HTMLParser):
    BLOCK = {"p","div","section","article","header","footer","address",
             "blockquote","pre","li","ul","ol","table","tr","h1","h2","h3","h4","h5","h6"}
    LINE = {"br","hr"}
    def __init__(self): super().__init__(); self.parts=[]
    def handle_starttag(self, t, a): t=t.lower(); self.parts.append("\n" if t in (self.LINE|self.BLOCK) else "")
    def handle_endtag(self, t):       t=t.lower(); self.parts.append("\n" if t in self.BLOCK else "")
    def handle_data(self, d):         self.parts.append(html.unescape(d) if d else "")
    def get_text(self):
        s="".join(self.parts)
        s=re.sub(r"\n{3,}", "\n\n", s)
        return "\n".join(line.rstrip() for line in s.splitlines()).strip()

def html_to_text(s: str) -> str:
    p=_HTMLToText(); p.feed(_to_str(s)); return p.get_text()
from PIL import Image
# ---------- OCR ----------
def ocr_image_with_confidence(img: Image.Image, lang: str) -> Tuple[str, None]:
    if not OCR_AVAILABLE:
        return "", None
    text = pytesseract.image_to_string(img, lang=lang)
    print("OCR")
    return text.strip(), None

def ocr_bytes(data: bytes, lang: str) -> Tuple[str, None]:
    if not OCR_AVAILABLE:
        return "", None
    try:
        img = Image.open(io.BytesIO(data))
        return ocr_image_with_confidence(img, lang)
    except Exception:
        return "", None

# ---------- PDF (FIXED OCR) ----------
def read_pdf(path: Path):
    start = time.time()
    if not PDF_AVAILABLE:
        raise RuntimeError("PyMuPDF not available")

    doc = fitz.open(str(path))
    text_parts = []
    ocr_pages_meta = []

    for i, page in enumerate(doc):
        text = page.get_text().strip()
        images = page.get_images(full=True)

        if not text and images:
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            ocr_text, _ = ocr_image_with_confidence(img, OCR_LANG)
            if ocr_text:
                text_parts.append(ocr_text)
                ocr_pages_meta.append({"page": i + 1})
        else:
            if text:
                text_parts.append(text)

    meta = {
        "type": "pdf",
        "engine": "PyMuPDF",
        "total_pages": doc.page_count,
        "total_ocr_pages": len(ocr_pages_meta),
        "ocr_pages": ocr_pages_meta,
        "total_time_sec": round(time.time() - start, 2),
    }
    return "\n\n".join(text_parts).strip(), meta

# ---------- DOCX ----------
from docx.oxml.ns import qn

def read_docx(path: Path):
    start = time.time()
    d = docx.Document(str(path))
    text_parts = []
    ocr_blocks_meta = []

    # Extract text from paragraphs
    for p in d.paragraphs:
        if p.text.strip():
            text_parts.append(p.text)

    # Extract text from tables
    for table in d.tables:
        for row in table.rows:
            text_parts.append("\t".join(c.text for c in row.cells))

    # A more robust way to find all images in the document
    image_parts = [
        rel.target_part
        for rel in d.part.rels.values()
        if rel.reltype.endswith('image')
    ]
    has_images = bool(image_parts)

    if has_images:
        if not OCR_AVAILABLE:
            logger.warning("OCR libraries not available, cannot extract text from images.")
        else:
            for i, image_part in enumerate(image_parts):
                try:
                    image_blob = image_part.blob
                    ocr_text, _ = ocr_bytes(image_blob, OCR_LANG)
                    if ocr_text:
                        text_parts.append(f"\n--- Text from image {i+1} ---\n{ocr_text}")
                        ocr_blocks_meta.append({"image_index": i + 1})
                except Exception as e:
                    logger.error(f"Error processing image {i+1}: {e}")

    meta = {
        "type": "docx",
        "has_images": has_images,
        "total_ocr_images": len(ocr_blocks_meta),
        "total_time_sec": round(time.time() - start, 2),
    }

    return "\n".join(text_parts).strip(), meta

# ---------- PPTX ----------
def _extract_from_shape(shape, slide_no, text_parts, ocr_blocks_meta, has_images):
    if hasattr(shape, "text") and shape.text.strip():
        text_parts.append(shape.text)

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        blob = shape.image.blob
        if blob:
            has_images[0] = True
            ocr_text, _ = ocr_bytes(blob, OCR_LANG)
            if ocr_text:
                text_parts.append(ocr_text)
                ocr_blocks_meta.append({"slide": slide_no})

    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            _extract_from_shape(s, slide_no, text_parts, ocr_blocks_meta, has_images)

def read_pptx(path: Path):
    start = time.time()
    prs = Presentation(str(path))
    text_parts = []
    ocr_blocks_meta = []
    has_images = [False]

    for slide_no, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            _extract_from_shape(shape, slide_no, text_parts, ocr_blocks_meta, has_images)

    meta = {
        "type": "pptx",
        "slides": len(prs.slides),
        "has_images": has_images[0],
        "total_ocr_images": len(ocr_blocks_meta),
        "total_time_sec": round(time.time() - start, 2),
    }
    return "\n".join(text_parts).strip(), meta

# ---------- OLD HANDLERS ----------
def read_txt(path: Path):
    return path.read_text(errors="ignore"), {"type": "txt"}

def read_csv(path: Path):
    if PANDAS_AVAILABLE:
        df = pd.read_csv(path, dtype=str).fillna("")
        return df.to_csv(index=False), {"type": "csv"}
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read(), {"type": "csv"}

def read_xlsx(path: Path):
    wb = load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            out.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(out), {"type": "xlsx"}

def read_email(path: Path):
    ext = path.suffix.lower()
    if ext == ".eml":
        if not EMAIL_AVAILABLE: raise RuntimeError("email parser unavailable")
        with open(path, "rb") as f:
            msg = BytesParser(policy=policy.default).parse(f)
        headers = (
            f"Subject: {_to_str(msg.get('subject',''))}\n"
            f"From: {_to_str(msg.get('from',''))}\n"
            f"To: {_to_str(msg.get('to',''))}\n"
            f"Date: {_to_str(msg.get('date',''))}\n\n"
        )
        plain, html_part = None, None
        parts = list(msg.walk()) if msg.is_multipart() else [msg]
        for part in parts:
            if part.get_content_disposition() == "attachment": continue
            ct = part.get_content_type()
            if ct == "text/plain" and plain is None: plain = _to_str(part.get_content())
            elif ct == "text/html" and html_part is None: html_part = _to_str(part.get_content())
        body = plain if plain is not None else (html_to_text(html_part) if html_part else "")
        return headers + body, {"type": "eml", "source": "text/plain" if plain else ("text/html->text" if html_part else "none")}

    elif ext == ".msg":
        if not MSG_AVAILABLE: raise RuntimeError("extract-msg not installed")
        m = extract_msg.Message(str(path))
        headers = (
            f"Subject: {_to_str(m.subject)}\n"
            f"From: {_to_str(m.sender)}\n"
            f"To: {_to_str(m.to)}\n"
            f"Date: {_to_str(m.date)}\n\n"
        )
        # Prefer HTML, then text, then RTF
        if getattr(m, "htmlBody", None):
            return headers + html_to_text(m.htmlBody), {"type": "msg", "source": "htmlBody"}
        if m.body:
            return headers + _to_str(m.body), {"type": "msg", "source": "body"}
        rtf = getattr(m, "rtf_decompressed", None) or getattr(m, "rtfBody", None)
        if rtf and (bt := rtf_to_text_quick(rtf)):
             return headers + bt, {"type": "msg", "source": "rtf"}
        # Fallback for attachments
        for att in (list(getattr(m, "attachments", []) or [])):
            name = (getattr(att, "longFilename", "") or getattr(att, "shortFilename", "") or "").lower()
            data = getattr(att, "data", b"")
            if name.endswith((".html",".htm")) and (txt := html_to_text(data).strip()):
                return headers + txt, {"type":"msg","source":f"attachment:{name}"}
            if name.endswith(".txt") and (txt := _to_str(data).strip()):
                return headers + txt, {"type":"msg","source":f"attachment:{name}"}
        return headers + "", {"type": "msg", "source": "none"}
    else:
        raise RuntimeError(f"Unsupported email type: {ext}")

# ---------- Dispatcher ----------
READERS = {
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".txt": read_txt,
    ".csv": read_csv,
    ".xlsx": read_xlsx,
    ".eml": read_email,
    ".msg": read_email,
}

def process_directory(input_dir: Path):
    for path in sorted(input_dir.rglob("*")):
        if path.is_file() and path.suffix.lower() in READERS:
            print("=" * 80)
            print(f"FILE: {path.name}")
            try:
                text, meta = READERS[path.suffix.lower()](path)
                print("[META]", meta)
                print("-" * 80)
                print(text)
            except Exception as e:
                logger.error(f"Failed to process {path.name}: {e}")
            print("=" * 80)

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", default="uploads")
    args = parser.parse_args()
    process_directory(Path(args.input))

from flask import Flask, request, jsonify
app = Flask(__name__)

@app.route("/extract", methods=["POST"])
def extract():
    try:
        uploaded_file = request.files.get("file")
        if not uploaded_file:
            return jsonify({"error": "No file uploaded"}), 400

        filename = uploaded_file.filename
        temp_path = Path(tempfile.gettempdir()) / filename
        uploaded_file.save(temp_path)

        reader = READERS.get(temp_path.suffix.lower())
        if not reader:
            return jsonify({"error": "Unsupported file type"}), 400

        text, meta = reader(temp_path)
        print(f"Extracted {len(text)} characters from {filename}")
        return jsonify({
            "filename": filename,
            "metadata": meta,
            "preview": text
        })
    except Exception as e:
        logger.exception("Extraction failed")
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1 and sys.argv[1] == "cli":
        main()
    else:
        app.run(host="0.0.0.0", port=5000)

